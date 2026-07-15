import json
from collections.abc import Callable
from pathlib import Path
from xml.sax.saxutils import escape
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.conversion import LocalConverter

FixtureWriter = Callable[[Path, str], None]


def write_text(path: Path, phrase: str) -> None:
    path.write_text(f"Heading\n\n{phrase}\n", encoding="utf-8")


def write_html(path: Path, phrase: str) -> None:
    path.write_text(
        f"<!doctype html><html><body><h1>Fixture</h1><p>{escape(phrase)}</p></body></html>",
        encoding="utf-8",
    )


def write_csv(path: Path, phrase: str) -> None:
    path.write_text(f"item,description\n1,{phrase}\n", encoding="utf-8")


def write_json(path: Path, phrase: str) -> None:
    path.write_text(json.dumps({"message": phrase}), encoding="utf-8")


def write_xml(path: Path, phrase: str) -> None:
    path.write_text(f"<?xml version='1.0'?><fixture>{escape(phrase)}</fixture>", encoding="utf-8")


def write_docx(path: Path, phrase: str) -> None:
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">\n'
        '  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.'
        'relationships+xml"/>\n'
        '  <Default Extension="xml" ContentType="application/xml"/>\n'
        '  <Override PartName="/word/document.xml" ContentType="application/vnd.'
        'openxmlformats-officedocument.wordprocessingml.document.main+xml"/>\n'
        "</Types>"
    )
    relationships = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n'
        '  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/'
        '2006/relationships/officeDocument" Target="word/document.xml"/>\n'
        "</Relationships>"
    )
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>{escape(phrase)}</w:t></w:r></w:p><w:sectPr/></w:body>
</w:document>"""
    with ZipFile(path, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)


def write_xlsx(path: Path, phrase: str) -> None:
    workbook = Workbook()
    workbook.active["A1"] = phrase
    workbook.save(path)


def write_pptx(path: Path, phrase: str) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    text_box.text_frame.text = phrase
    presentation.save(path)


def write_pdf(path: Path, phrase: str) -> None:
    safe_phrase = phrase.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 18 Tf 72 720 Td ({safe_phrase}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode())
        content.extend(body)
        content.extend(b"\nendobj\n")
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    trailer = (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n"
    )
    content.extend(trailer.encode())
    path.write_bytes(content)


FORMATS: list[tuple[str, FixtureWriter, str]] = [
    (".txt", write_text, "Plain text matrix phrase 1042"),
    (".html", write_html, "HTML matrix phrase 2153"),
    (".csv", write_csv, "CSV matrix phrase 3264"),
    (".json", write_json, "JSON matrix phrase 4375"),
    (".xml", write_xml, "XML matrix phrase 5486"),
    (".docx", write_docx, "Word matrix phrase 6597"),
    (".xlsx", write_xlsx, "Excel matrix phrase 7608"),
    (".pptx", write_pptx, "PowerPoint matrix phrase 8719"),
    (".pdf", write_pdf, "PDF matrix phrase 9820"),
]


@pytest.mark.parametrize(("suffix", "writer", "phrase"), FORMATS, ids=[item[0] for item in FORMATS])
def test_real_markitdown_format_conversion(
    tmp_path: Path,
    suffix: str,
    writer: FixtureWriter,
    phrase: str,
) -> None:
    source = tmp_path / f"source{suffix}"
    output = tmp_path / f"source{suffix}.md"
    writer(source, phrase)

    result = LocalConverter(timeout_seconds=60).convert(source, output)

    assert result.outcome == "converted", result.reason
    assert output.suffix == ".md"
    assert phrase in output.read_text(encoding="utf-8")
